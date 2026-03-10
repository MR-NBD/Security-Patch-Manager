#!/usr/bin/env python3
"""AI per le IT Operations — genera presentazione PowerPoint.

Uso:  python3 gen_presentation.py
Output: AI_IT_Operations_n8n.pptx nella stessa cartella.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE

# ─── Palette ──────────────────────────────────────────────────────────────────
BG     = RGBColor( 10,  22,  40)   # navy profondo — sfondo base
CARD   = RGBColor( 18,  38,  66)   # card
CARD2  = RGBColor( 24,  50,  86)   # card alternativa
BLUE   = RGBColor( 74, 158, 255)   # accent primario
TEAL   = RGBColor(  0, 212, 170)   # accent secondario
ORANGE = RGBColor(255, 107,  53)   # highlight
RED    = RGBColor(232,  72,  85)   # critico / cons
GREEN  = RGBColor(  0, 186, 136)   # success / pros
WHITE  = RGBColor(255, 255, 255)
LGRAY  = RGBColor(148, 170, 196)   # testo secondario
YELLOW = RGBColor(255, 213,  79)   # novelty

W = Inches(13.33)   # 16:9
H = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank


def fill_bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or BG


def R(slide, x, y, w, h, fill, rx=False, border_color=None):
    """Rettangolo o rettangolo arrotondato."""
    stype = SHAPE.ROUNDED_RECTANGLE if rx else SHAPE.RECTANGLE
    s = slide.shapes.add_shape(stype, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def C(slide, cx, cy, d, fill):
    """Cerchio centrato in (cx, cy) con diametro d."""
    s = slide.shapes.add_shape(SHAPE.OVAL, cx - d / 2, cy - d / 2, d, d)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def ARR(slide, x, y, w, h, fill):
    """Freccia destra."""
    s = slide.shapes.add_shape(SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def T(slide, text, x, y, w, h, size, color,
      bold=False, align=PP_ALIGN.LEFT, italic=False):
    """Textbox singolo paragrafo."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size    = Pt(size)
    r.font.color.rgb = color
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.name    = "Calibri"
    return tb


def hint(slide, text):
    """Striscia in basso con suggerimento visuale per il presentatore."""
    R(slide, Inches(0.5), H - Inches(0.45), Inches(12.3), Inches(0.36),
      RGBColor(5, 25, 45), rx=True)
    T(slide, f"[  {text}  ]",
      Inches(0.7), H - Inches(0.43), Inches(12.0), Inches(0.32),
      10, RGBColor(60, 100, 130), italic=True, align=PP_ALIGN.CENTER)


# ─── Slide builders ───────────────────────────────────────────────────────────

def slide_01_title(prs):
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.1), BLUE)
    R(sl, 0, H - Inches(0.1), W, Inches(0.1), TEAL)

    T(sl, "AI per le IT Operations",
      Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.9),
      54, WHITE, bold=True, align=PP_ALIGN.CENTER)

    T(sl, "Automazione · Intelligenza · Innovazione",
      Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.75),
      28, TEAL, align=PP_ALIGN.CENTER)

    T(sl, "n8n  ×  UYUNI  ×  SaltStack  ×  AI Open Source",
      Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.6),
      18, LGRAY, align=PP_ALIGN.CENTER)

    # accent bar centrale
    R(sl, Inches(4.5), Inches(5.25), Inches(4.3), Inches(0.05), BLUE)

    T(sl, "[ inserire logo aziendale ]",
      Inches(0.5), H - Inches(0.7), Inches(4), Inches(0.45),
      11, RGBColor(50, 70, 90), italic=True)
    return sl


def slide_02_benefits(prs):
    """5 blocchi benefici — stesso stile della reference."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), BLUE)

    T(sl, "Perché l'AI per le IT Operations",
      Inches(0.8), Inches(0.18), Inches(11.7), Inches(0.7),
      28, WHITE, bold=True, align=PP_ALIGN.CENTER)

    items = [
        ("⚙",  "Automatizzare\ntask ripetitivi\ne noiosi",       BLUE),
        ("📊", "Analizzare\ngrandi volumi\ndi dati",             TEAL),
        ("🔗", "Correlare info\ntra sistemi\neterogenei",        ORANGE),
        ("🚀", "Potenziare\nle capacità\ndegli operatori",       GREEN),
        ("⚖",  "Ridurre\nil carico sul\nteam Operations",        YELLOW),
    ]
    cw, ch = Inches(2.3), Inches(5.2)
    gap    = Inches(0.25)
    sx     = Inches(0.55)
    y      = Inches(1.2)

    for i, (icon, text, color) in enumerate(items):
        x = sx + i * (cw + gap)
        R(sl, x, y, cw, ch, CARD, rx=True)
        R(sl, x, y, cw, Inches(0.12), color)
        # cerchio icona
        C(sl, x + cw / 2, y + Inches(1.05), Inches(0.95), color)
        T(sl, icon,
          x + cw / 2 - Inches(0.38), y + Inches(0.67),
          Inches(0.76), Inches(0.76), 24, WHITE, align=PP_ALIGN.CENTER)
        T(sl, text,
          x + Inches(0.15), y + Inches(1.85), cw - Inches(0.3), Inches(2.8),
          15, WHITE, align=PP_ALIGN.CENTER)
    return sl


def slide_03_reference_scenario(prs):
    """Sorgenti dati — scenario di riferimento."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), BLUE)

    T(sl, "Scenario di Riferimento — Le Sorgenti Dati Ops",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      24, WHITE, bold=True)

    sources = [
        ("VM On-Premise",     "Linux · Windows · bare metal",  BLUE),
        ("Azure / AWS / GCP", "Cloud multi-provider",           TEAL),
        ("CMDB",              "Asset inventory & topology",     ORANGE),
        ("Log & Metriche",    "Traces · Events · Alerts",       GREEN),
        ("Errata / CVE",      "Security advisories NVD/vendor", RED),
        ("Service Desk",      "Ticket & incident management",   YELLOW),
    ]
    bw, bh = Inches(3.7), Inches(1.45)
    gx, gy = Inches(0.5), Inches(0.35)
    sx, sy = Inches(0.5), Inches(1.2)

    for i, (title, sub, color) in enumerate(sources):
        col, row = i % 3, i // 3
        x = sx + col * (bw + gx)
        y = sy + row * (bh + gy)
        R(sl, x, y, bw, bh, CARD, rx=True)
        R(sl, x, y, Inches(0.1), bh, color)
        T(sl, title, x + Inches(0.25), y + Inches(0.15), bw - Inches(0.3), Inches(0.55),
          17, WHITE, bold=True)
        T(sl, sub,   x + Inches(0.25), y + Inches(0.75), bw - Inches(0.3), Inches(0.5),
          13, LGRAY)

    # freccia → ops team
    ARR(sl, Inches(11.6), Inches(2.8), Inches(1.1), Inches(0.45), BLUE)
    R(sl, Inches(12.0), Inches(1.9), Inches(1.1), Inches(2.1), CARD2, rx=True)
    T(sl, "OPS\nTEAM", Inches(12.05), Inches(2.4), Inches(1.0), Inches(1.1),
      14, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # problema in basso
    R(sl, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.72),
      RGBColor(38, 18, 18), rx=True)
    T(sl, "⚠  Eterogeneità delle fonti  →  nessuna visione unificata  →  team costantemente sovraccarico",
      Inches(0.8), Inches(6.58), Inches(12.0), Inches(0.5),
      14, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    return sl


def slide_04_how_we_got_here(prs):
    """Timeline: infrastruttura → patch management → AI."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), ORANGE)

    T(sl, "Come Siamo Arrivati Qui",
      Inches(0.6), Inches(0.22), Inches(12), Inches(0.68),
      30, WHITE, bold=True)

    steps = [
        ("①", "Infrastruttura\nCentralizzata",
         "UYUNI + SaltStack per\ngestione multi-OS e multi-tenant\nCSP-independent", BLUE),
        ("②", "Patch Management\n(Progetto Parallelo)",
         "Security Patch Manager:\nla necessità di automazione\nha reso visibile il potenziale", ORANGE),
        ("③", "AI per le\nIT Operations",
         "L'infrastruttura esistente\ndiventa la base per\nl'integrazione con n8n + AI", TEAL),
    ]
    sw, sh = Inches(3.6), Inches(4.0)
    gap    = Inches(0.55)
    sx, y  = Inches(0.55), Inches(1.6)

    for i, (num, title, desc, color) in enumerate(steps):
        x = sx + i * (sw + gap)
        R(sl, x, y, sw, sh, CARD, rx=True)
        R(sl, x, y, sw, Inches(0.12), color)
        C(sl, x + sw / 2, y + Inches(0.95), Inches(0.85), color)
        T(sl, num,
          x + sw / 2 - Inches(0.33), y + Inches(0.62),
          Inches(0.66), Inches(0.66), 22, WHITE, bold=True, align=PP_ALIGN.CENTER)
        T(sl, title,
          x, y + Inches(1.65), sw, Inches(0.85),
          16, WHITE, bold=True, align=PP_ALIGN.CENTER)
        T(sl, desc,
          x + Inches(0.2), y + Inches(2.6), sw - Inches(0.4), Inches(1.15),
          13, LGRAY, align=PP_ALIGN.CENTER)
        # freccia tra step
        if i < 2:
            ARR(sl, x + sw + Inches(0.1), y + sh / 2 - Inches(0.22),
                Inches(0.35), Inches(0.38), color)

    T(sl, "L'infrastruttura costruita per il patch management è oggi il fondamento dell'AI Operations",
      Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.6),
      16, TEAL, bold=True, align=PP_ALIGN.CENTER)
    return sl


def slide_05_uyuni(prs):
    """UYUNI come strumento infrastrutturale."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), TEAL)

    # Sinistra: title + descrizione
    T(sl, "UYUNI",
      Inches(0.6), Inches(0.7), Inches(5.8), Inches(1.3),
      52, WHITE, bold=True)
    T(sl, "Gestione Centralizzata\ndell'Infrastruttura",
      Inches(0.6), Inches(2.0), Inches(5.8), Inches(0.9),
      20, TEAL)
    T(sl, ("Strumento open source per la gestione unificata\n"
           "di sistemi Linux eterogenei.\n\n"
           "Registra, monitora e governa fleet multi-OS\n"
           "da un singolo pannello di controllo.\n\n"
           "Usato come infrastruttura di base su cui\n"
           "n8n si aggancia per l'AI orchestration."),
      Inches(0.6), Inches(3.0), Inches(5.8), Inches(2.8),
      13, LGRAY)

    # XML-RPC highlight
    R(sl, Inches(0.6), Inches(5.9), Inches(5.5), Inches(0.85), CARD2, rx=True)
    R(sl, Inches(0.6), Inches(5.9), Inches(0.1), Inches(0.85), ORANGE)
    T(sl, "XML-RPC API  —  Il punto di aggancio per n8n e l'AI layer",
      Inches(0.85), Inches(6.05), Inches(5.1), Inches(0.5),
      13, ORANGE, bold=True)

    # Destra: 4 feature card
    features = [
        ("🖥", "Multi-OS",        "Ubuntu · RHEL · SUSE · Debian",              BLUE),
        ("📦", "Errata & Patch",  "Distribuzione centralizzata di aggiornamenti",TEAL),
        ("⚙",  "SaltStack",       "Configuration management integrato nativo",    ORANGE),
        ("🔌", "XML-RPC API",     "Integrazione con n8n, AI tools, automazioni",  GREEN),
    ]
    fw, fh = Inches(6.7), Inches(1.45)
    gap    = Inches(0.22)
    fx, fy = Inches(6.4), Inches(0.85)

    for i, (icon, title, desc, color) in enumerate(features):
        y = fy + i * (fh + gap)
        R(sl, fx, y, fw, fh, CARD, rx=True)
        R(sl, fx, y, Inches(0.09), fh, color)
        C(sl, fx + Inches(0.68), y + fh / 2, Inches(0.58), color)
        T(sl, icon, fx + Inches(0.41), y + Inches(0.38), Inches(0.58), Inches(0.58),
          16, WHITE, align=PP_ALIGN.CENTER)
        T(sl, title, fx + Inches(1.05), y + Inches(0.12), fw - Inches(1.1), Inches(0.48),
          15, WHITE, bold=True)
        T(sl, desc,  fx + Inches(1.05), y + Inches(0.65), fw - Inches(1.1), Inches(0.52),
          12, LGRAY)

    hint(sl, "SUGGERIMENTO: inserire screenshot dashboard UYUNI con lista sistemi registrati")
    return sl


def slide_06_oss_numbers(prs):
    """OSS Numbers — UYUNI, SaltStack."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), TEAL)

    T(sl, "Open Source — La Solidità dei Numeri",
      Inches(0.6), Inches(0.22), Inches(12), Inches(0.65),
      28, WHITE, bold=True, align=PP_ALIGN.CENTER)

    stats = [
        ("15+",  "anni di sviluppo",    "UYUNI / Spacewalk",         BLUE),
        ("800+", "stelle GitHub",        "UYUNI",                     TEAL),
        ("14k+", "stelle GitHub",        "SaltStack",                 ORANGE),
        ("13+",  "anni in produzione",   "SaltStack",                 GREEN),
        ("100%", "CSP-Independent",      "Nessun vendor lock-in",     YELLOW),
    ]
    cw, ch = Inches(2.3), Inches(4.5)
    gap    = Inches(0.28)
    sx, y  = Inches(0.55), Inches(1.4)

    for i, (number, label, tool, color) in enumerate(stats):
        x = sx + i * (cw + gap)
        R(sl, x, y, cw, ch, CARD, rx=True)
        R(sl, x, y, cw, Inches(0.12), color)
        T(sl, number,
          x, y + Inches(0.75), cw, Inches(1.1),
          40, color, bold=True, align=PP_ALIGN.CENTER)
        T(sl, label,
          x + Inches(0.1), y + Inches(1.9), cw - Inches(0.2), Inches(0.75),
          14, WHITE, align=PP_ALIGN.CENTER)
        T(sl, tool,
          x + Inches(0.1), y + Inches(2.75), cw - Inches(0.2), Inches(0.65),
          12, LGRAY, align=PP_ALIGN.CENTER)

    R(sl, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.65), CARD2, rx=True)
    T(sl, "Foundation: openSUSE · VMware Salt · Community globale attiva · Codice auditabile · Costi zero di licenza",
      Inches(0.8), Inches(6.22), Inches(12.0), Inches(0.45),
      13, LGRAY, align=PP_ALIGN.CENTER)
    return sl


def slide_07_section_pillars(prs):
    """Divisore sezione — I Tre Pilastri."""
    sl = new_slide(prs)
    fill_bg(sl, BLUE)
    R(sl, 0, 0, W, Inches(0.1), WHITE)
    R(sl, 0, H - Inches(0.1), W, Inches(0.1), WHITE)

    T(sl, "I Tre Pilastri",
      Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.7),
      60, WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(sl, "del Valore dell'AI",
      Inches(1.0), Inches(3.4), Inches(11.3), Inches(1.1),
      38, RGBColor(180, 215, 255), align=PP_ALIGN.CENTER)
    T(sl, "Automazione  ·  Intelligence  ·  Novelty",
      Inches(1.0), Inches(4.75), Inches(11.3), Inches(0.7),
      22, RGBColor(220, 235, 255), align=PP_ALIGN.CENTER)
    return sl


def slide_08_three_pillars(prs):
    """3 pilastri contenuto."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), BLUE)

    T(sl, "I Tre Pilastri del Valore dell'AI per le Operations",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      22, WHITE, bold=True, align=PP_ALIGN.CENTER)

    pillars = [
        ("AUTOMAZIONE", BLUE,   "⚙",
         ["Elimina il lavoro manuale",
          "ripetitivo e time-consuming",
          "",
          "→ Classificazione alert",
          "→ Scheduling intelligente",
          "→ Report automatici notturni",
          "→ Ticket enrichment"]),
        ("INTELLIGENCE", TEAL,  "🧠",
         ["Analizza ciò che l'uomo",
          "non può processare da solo",
          "",
          "→ Correlazione log/metriche",
          "→ Root cause analysis",
          "→ Anomaly detection su KPI",
          "→ Pattern recognition"]),
        ("NOVELTY",      YELLOW, "✨",
         ["Ricerca proattiva",
          "di situazioni nuove",
          "",
          "→ Anomalie prima dell'incidente",
          "→ Suggerimento azioni inedite",
          "→ Knowledge base evolutiva",
          "→ Proactive alerting"]),
    ]
    cw, ch = Inches(3.95), Inches(6.0)
    gap    = Inches(0.35)
    sx, y  = Inches(0.55), Inches(1.0)

    for i, (title, color, icon, pts) in enumerate(pillars):
        x = sx + i * (cw + gap)
        R(sl, x, y, cw, ch, CARD, rx=True)
        R(sl, x, y, cw, Inches(0.15), color)
        C(sl, x + cw / 2, y + Inches(1.0), Inches(1.05), color)
        T(sl, icon,
          x + cw / 2 - Inches(0.42), y + Inches(0.57),
          Inches(0.84), Inches(0.84), 26, WHITE, align=PP_ALIGN.CENTER)
        T(sl, title,
          x, y + Inches(1.72), cw, Inches(0.55),
          15, color, bold=True, align=PP_ALIGN.CENTER)
        for j, pt in enumerate(pts):
            py = y + Inches(2.42) + j * Inches(0.46)
            clr = WHITE if pt.startswith("→") else LGRAY
            T(sl, pt, x + Inches(0.22), py, cw - Inches(0.28), Inches(0.43),
              13 if pt.startswith("→") else 13, clr)
    return sl


def slide_09_real_scenarios(prs):
    """Scenari reali — 3 colonne (stesso stile della reference)."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), BLUE)

    T(sl, "Scenari Reali",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      28, WHITE, bold=True, align=PP_ALIGN.CENTER)

    cols = [
        ("RETRIEVAL", BLUE, "🔍",
         ["Ricerca documentazione infrastruttura",
          "e deployment",
          "Correlazione metadati server e CVE",
          "Best Practice Assessment",
          "(SQL Server, Linux Server)",
          "Inventory e change tracking"]),
        ("PROCESS AUTOMATION", TEAL, "🔄",
         ["Morning Check automatico:",
          "stato piattaforma ogni mattina",
          "Automated Report su checklist",
          "Ticket enrichment via",
          "check automatici",
          "Smart alerting con pre-analisi"]),
        ("REACTIVE / PROACTIVE", ORANGE, "⚡",
         ["Verifica qualità del servizio",
          "per applicazione/sistema",
          "Rilevamento anomalie su KPI",
          "Root cause analysis post-incidente",
          "Build evidence prima",
          "di un war room"]),
    ]
    cw, ch = Inches(3.95), Inches(5.9)
    gap    = Inches(0.35)
    sx, y  = Inches(0.55), Inches(1.05)

    for i, (title, color, icon, items) in enumerate(cols):
        x = sx + i * (cw + gap)
        R(sl, x, y, cw, ch, CARD, rx=True)
        R(sl, x, y, cw, Inches(0.15), color)
        C(sl, x + cw / 2, y + Inches(0.95), Inches(0.9), color)
        T(sl, icon,
          x + cw / 2 - Inches(0.36), y + Inches(0.58),
          Inches(0.72), Inches(0.72), 22, WHITE, align=PP_ALIGN.CENTER)
        T(sl, title,
          x, y + Inches(1.58), cw, Inches(0.52),
          13, color, bold=True, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            iy = y + Inches(2.2) + j * Inches(0.55)
            bullet = "·  " if j % 2 == 0 else "    "
            T(sl, bullet + item, x + Inches(0.18), iy, cw - Inches(0.25), Inches(0.52),
              12, WHITE if j % 2 == 0 else LGRAY)
    return sl


def slide_10_orchestration(prs):
    """AI Orchestration Layer — diagramma a strati."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), BLUE)

    T(sl, "AI Orchestration Layer — n8n al Centro",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      26, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Layer 3 — Output (operatore)
    R(sl, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.1), CARD, rx=True)
    R(sl, Inches(0.5), Inches(1.0), Inches(0.1), Inches(1.1), TEAL)
    T(sl, "OPERATORE  /  AI MODELS  /  ALERT  /  REPORT  /  DASHBOARD",
      Inches(0.75), Inches(1.2), Inches(11.8), Inches(0.65),
      18, TEAL, bold=True, align=PP_ALIGN.CENTER)

    T(sl, "▲  OUTPUT  ▲",
      Inches(5.8), Inches(2.2), Inches(1.75), Inches(0.4),
      12, LGRAY, align=PP_ALIGN.CENTER)

    # Layer 2 — n8n
    R(sl, Inches(1.8), Inches(2.7), Inches(9.7), Inches(1.45),
      RGBColor(0, 60, 50), rx=True, border_color=TEAL)
    T(sl, "n8n  —  AI Orchestration & Workflow Layer",
      Inches(1.8), Inches(2.92), Inches(9.7), Inches(0.65),
      20, TEAL, bold=True, align=PP_ALIGN.CENTER)
    T(sl, "Workflow visuale  ·  Modelli AI  ·  Trasparenza  ·  Audit Trail  ·  Trigger event-driven",
      Inches(1.8), Inches(3.6), Inches(9.7), Inches(0.4),
      12, LGRAY, align=PP_ALIGN.CENTER)

    T(sl, "▼  QUERY / AZIONI  ▼",
      Inches(5.6), Inches(4.25), Inches(2.1), Inches(0.4),
      12, LGRAY, align=PP_ALIGN.CENTER)

    # Layer 1 — Infrastructure
    srcs = [
        ("UYUNI\nXML-RPC",  BLUE),
        ("Log &\nMetriche",  TEAL),
        ("CMDB",             ORANGE),
        ("Errata\n& CVE",    RED),
        ("Service\nDesk",    GREEN),
    ]
    bw, bh = Inches(2.15), Inches(1.4)
    gap    = Inches(0.24)
    sx, sy = Inches(0.5), Inches(4.8)

    for i, (name, color) in enumerate(srcs):
        x = sx + i * (bw + gap)
        R(sl, x, sy, bw, bh, CARD, rx=True)
        R(sl, x, sy, bw, Inches(0.1), color)
        T(sl, name, x, sy + Inches(0.3), bw, Inches(0.85),
          14, WHITE, bold=True, align=PP_ALIGN.CENTER)

    R(sl, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.4), CARD2, rx=True)
    T(sl, "INFRASTRUTTURA — SORGENTI DATI",
      Inches(0.5), Inches(6.43), Inches(12.3), Inches(0.28),
      12, LGRAY, align=PP_ALIGN.CENTER)

    hint(sl, "SUGGERIMENTO: sostituire con diagramma animato (draw.io / Miro) mostrando i flussi dati")
    return sl


def slide_11_section_challenges(prs):
    """Divisore sezione — Le Sfide."""
    sl = new_slide(prs)
    fill_bg(sl, RGBColor(38, 10, 15))
    R(sl, 0, 0, W, Inches(0.1), RED)
    R(sl, 0, H - Inches(0.1), W, Inches(0.1), RED)

    T(sl, "Le Sfide dell'AI",
      Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.9),
      60, WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(sl, "Quando l'intelligenza non spiega sé stessa",
      Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.0),
      26, RGBColor(255, 155, 165), align=PP_ALIGN.CENTER)
    T(sl, "Trasparenza  ·  Interpretabilità  ·  Audit Trail",
      Inches(1.0), Inches(4.75), Inches(11.3), Inches(0.65),
      18, RED, align=PP_ALIGN.CENTER)
    return sl


def slide_12_challenges(prs):
    """Criticità: Trasparenza e Interpretabilità."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), RED)

    T(sl, "Le Criticità dell'AI in Produzione",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      26, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Card sinistra — Trasparenza
    R(sl, Inches(0.5), Inches(1.05), Inches(5.9), Inches(5.6), CARD, rx=True)
    R(sl, Inches(0.5), Inches(1.05), Inches(5.9), Inches(0.15), RED)
    C(sl, Inches(3.45), Inches(2.3), Inches(1.1), RGBColor(75, 18, 24))
    T(sl, "🚫",
      Inches(2.95), Inches(1.88), Inches(1.0), Inches(0.82),
      28, RED, align=PP_ALIGN.CENTER)
    T(sl, "CARENZA DI\nTRASPARENZA",
      Inches(0.7), Inches(2.95), Inches(5.5), Inches(0.95),
      20, RED, bold=True, align=PP_ALIGN.CENTER)
    T(sl, "Il modello decide\nma non mostra il ragionamento",
      Inches(0.7), Inches(3.95), Inches(5.5), Inches(0.7),
      14, WHITE, align=PP_ALIGN.CENTER)
    T(sl, ("→ Impossibile auditare una decisione\n"
           "→ Nessun trail: perché ha fatto X?\n"
           "→ Compliance: come dimostro la logica?"),
      Inches(0.85), Inches(4.75), Inches(5.4), Inches(1.5),
      13, LGRAY)

    # Card destra — Interpretabilità
    R(sl, Inches(6.9), Inches(1.05), Inches(5.9), Inches(5.6), CARD, rx=True)
    R(sl, Inches(6.9), Inches(1.05), Inches(5.9), Inches(0.15), ORANGE)
    C(sl, Inches(9.85), Inches(2.3), Inches(1.1), RGBColor(75, 38, 10))
    T(sl, "❓",
      Inches(9.35), Inches(1.88), Inches(1.0), Inches(0.82),
      28, ORANGE, align=PP_ALIGN.CENTER)
    T(sl, "CARENZA DI\nINTERPRETABILITÀ",
      Inches(7.1), Inches(2.95), Inches(5.5), Inches(0.95),
      20, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    T(sl, "Non si capisce perché\nuna decisione è stata presa",
      Inches(7.1), Inches(3.95), Inches(5.5), Inches(0.7),
      14, WHITE, align=PP_ALIGN.CENTER)
    T(sl, ("→ Debug impossibile su incidenti\n"
           "→ Fiducia zero dall'operatore\n"
           "→ Rischio azioni errate in produzione"),
      Inches(7.1), Inches(4.75), Inches(5.4), Inches(1.5),
      13, LGRAY)

    # Banner bottom
    R(sl, Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.52),
      RGBColor(48, 18, 18), rx=True)
    T(sl, "Serve uno strato che renda ogni azione AI:  VISIBILE  ·  TRACCIABILE  ·  COMPRENSIBILE",
      Inches(0.8), Inches(6.9), Inches(12.0), Inches(0.38),
      15, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    return sl


def slide_13_section_n8n(prs):
    """Divisore sezione — n8n."""
    sl = new_slide(prs)
    fill_bg(sl, RGBColor(0, 46, 38))
    R(sl, 0, 0, W, Inches(0.1), TEAL)
    R(sl, 0, H - Inches(0.1), W, Inches(0.1), TEAL)

    T(sl, "n8n",
      Inches(1.0), Inches(1.2), Inches(11.3), Inches(2.1),
      82, WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(sl, "La Risposta Ingegneristica",
      Inches(1.0), Inches(3.3), Inches(11.3), Inches(1.0),
      34, TEAL, align=PP_ALIGN.CENTER)
    T(sl, "Workflow Visuale  ·  Trasparenza  ·  Modelli Locali",
      Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.65),
      18, RGBColor(175, 238, 225), align=PP_ALIGN.CENTER)
    T(sl, "50.000+ ★ GitHub  ·  400+ integrazioni  ·  80.000+ community  ·  Open Source",
      Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.55),
      15, RGBColor(110, 195, 180), align=PP_ALIGN.CENTER)
    return sl


def slide_14_n8n_transparency(prs):
    """n8n — Trasparenza Nativa."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), TEAL)

    T(sl, "n8n — Workflow Visuale = Trasparenza Nativa",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      24, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Workflow nodes visivi (top)
    node_colors  = [BLUE, TEAL, ORANGE, GREEN]
    node_labels  = ["TRIGGER\n(UYUNI Event)", "AI NODE\n(Analisi CVE)",
                    "DECISION\n(Risk Score)", "ACTION\n(Notify/Apply)"]
    nw, nh = Inches(2.55), Inches(1.1)
    nsx, ny = Inches(0.4), Inches(1.0)

    for i, (color, label) in enumerate(zip(node_colors, node_labels)):
        x = nsx + i * Inches(3.1)
        R(sl, x, ny, nw, nh, CARD, rx=True)
        R(sl, x, ny, nw, Inches(0.1), color)
        T(sl, label, x, ny + Inches(0.18), nw, nh - Inches(0.22),
          12, WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            ARR(sl, x + nw + Inches(0.08), ny + nh / 2 - Inches(0.18),
                Inches(0.45), Inches(0.34), color)

    # Log strip
    R(sl, Inches(0.4), Inches(2.25), Inches(12.5), Inches(0.48), CARD2, rx=True)
    T(sl, "📋  Execution Log: ogni nodo registra  input · output · tempo · modello · errori",
      Inches(0.65), Inches(2.3), Inches(12.2), Inches(0.4),
      13, TEAL, bold=True)

    # 3 benefit cards
    benefits = [
        (BLUE,   "AUDIT TRAIL AUTOMATICO",
                 "Ogni step del flusso è loggato.\nSai esattamente cosa è entrato,\ncosa è uscito e perché."),
        (TEAL,   "DEBUG NODO PER NODO",
                 "Non si debugga un monolite.\nIsoli il problema nel singolo\nnodo del workflow."),
        (ORANGE, "WORKFLOW ISOLATI PER TENANT",
                 "I dati di ogni cliente restano\nseparati nei propri workflow\ndedicati."),
    ]
    bw, bh = Inches(3.9), Inches(1.75)
    bsx, by = Inches(0.5), Inches(2.92)

    for i, (color, title, desc) in enumerate(benefits):
        x = bsx + i * (bw + Inches(0.45))
        R(sl, x, by, bw, bh, CARD, rx=True)
        R(sl, x, by, Inches(0.09), bh, color)
        T(sl, title, x + Inches(0.22), by + Inches(0.12), bw - Inches(0.28), Inches(0.48),
          13, color, bold=True)
        T(sl, desc,  x + Inches(0.22), by + Inches(0.65), bw - Inches(0.28), Inches(0.92),
          12, LGRAY)

    hint(sl, "SUGGERIMENTO: inserire screenshot reale di workflow n8n con nodi colorati + execution log")
    return sl


def slide_15_n8n_engineering(prs):
    """n8n — Ingegneria Personalizzata dei Modelli."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), TEAL)

    T(sl, "n8n — Ingegneria Personalizzata dei Modelli AI",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      24, WHITE, bold=True, align=PP_ALIGN.CENTER)

    cards = [
        (BLUE,   "🎯  SELEZIONE MODELLO PER USE CASE",
                 ("Modello leggero per classificare alert.\n"
                  "Modello potente per analisi CVE complessa.\n"
                  "Ogni task ha il modello giusto, senza sprechi.")),
        (TEAL,   "🔀  SEGMENTAZIONE DEI FLUSSI",
                 ("Discovery · Test · Approvazione · Notifica.\n"
                  "Ogni processo: workflow dedicato.\n"
                  "Manutenibilità e isolamento totale.")),
        (ORANGE, "🔍  DEBUG GRANULARE",
                 ("Il problema è nel nodo X, non nel sistema.\n"
                  "Test isolato per segmento.\n"
                  "Fine-tuning senza toccare il resto.")),
        (GREEN,  "🧬  SPECIALIZZAZIONE & FINE-TUNING",
                 ("Modelli addestrati su domini specifici:\n"
                  "Security · Compliance · Infrastruttura.\n"
                  "Knowledge base proprietaria e aggiornabile.")),
    ]
    cw, ch = Inches(5.95), Inches(2.65)
    gap    = Inches(0.35)
    sx, sy = Inches(0.55), Inches(1.05)

    for idx, (color, title, desc) in enumerate(cards):
        ci, ri = idx % 2, idx // 2
        x = sx + ci * (cw + gap)
        y = sy + ri * (ch + gap)
        R(sl, x, y, cw, ch, CARD, rx=True)
        R(sl, x, y, cw, Inches(0.12), color)
        T(sl, title, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.3), Inches(0.85),
          16, color, bold=True)
        T(sl, desc,  x + Inches(0.25), y + Inches(1.12), cw - Inches(0.3), Inches(1.35),
          13, LGRAY)
    return sl


def slide_16_local_models(prs):
    """Sicurezza — Modelli in Rete Locale."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), GREEN)

    T(sl, "Sicurezza — I Dati dei Clienti Non Escono Mai",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      26, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Box rete locale
    R(sl, Inches(0.5), Inches(1.05), Inches(8.0), Inches(5.75),
      RGBColor(7, 28, 20), rx=True, border_color=GREEN)
    T(sl, "🔒  RETE LOCALE / ON-PREMISE",
      Inches(0.72), Inches(1.15), Inches(7.6), Inches(0.5),
      14, GREEN, bold=True)

    inner = [
        ("n8n\nOrchestrator",  TEAL,   Inches(0.8),  Inches(1.85)),
        ("Ollama\n(LLM locale)", GREEN, Inches(3.1),  Inches(1.85)),
        ("UYUNI\nXML-RPC",     BLUE,   Inches(5.4),  Inches(1.85)),
        ("PostgreSQL\nDB",      ORANGE, Inches(0.8),  Inches(3.65)),
        ("Dati\nClienti",       YELLOW, Inches(3.1),  Inches(3.65)),
        ("SPM\nOrchestrator",   LGRAY,  Inches(5.4),  Inches(3.65)),
    ]
    for label, color, ix, iy in inner:
        R(sl, ix, iy, Inches(2.0), Inches(1.35), CARD, rx=True)
        R(sl, ix, iy, Inches(2.0), Inches(0.1), color)
        T(sl, label, ix, iy + Inches(0.22), Inches(2.0), Inches(0.95),
          13, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Destra: 4 benefit
    ritems = [
        (GREEN,  "GDPR COMPLIANT",
                 "Nessun dato cliente inviato\na servizi cloud pubblici"),
        (TEAL,   "ZERO VENDOR LOCK-IN AI",
                 "Ollama, vLLM, LM Studio:\nopen source self-hosted"),
        (BLUE,   "LATENZA RIDOTTA",
                 "Inference locale = risposta\npiù rapida, costi contenuti"),
        (ORANGE, "CONTROLLO TOTALE",
                 "Versioning modelli, fine-tuning\ninterno, audit completo"),
    ]
    rx0, ry0 = Inches(8.8), Inches(1.05)
    rw, rh   = Inches(4.2), Inches(1.25)
    rgap     = Inches(0.22)

    for i, (color, title, desc) in enumerate(ritems):
        ry = ry0 + i * (rh + rgap)
        R(sl, rx0, ry, rw, rh, CARD, rx=True)
        R(sl, rx0, ry, Inches(0.08), rh, color)
        T(sl, title, rx0 + Inches(0.2), ry + Inches(0.1), rw - Inches(0.25), Inches(0.48),
          13, color, bold=True)
        T(sl, desc,  rx0 + Inches(0.2), ry + Inches(0.6), rw - Inches(0.25), Inches(0.55),
          12, LGRAY)

    # Cloud bloccato
    R(sl, Inches(8.8), Inches(6.28), Inches(4.2), Inches(0.58),
      RGBColor(48, 18, 18), rx=True)
    T(sl, "✗  Cloud AI pubblici: traffico clienti non esposto",
      Inches(9.0), Inches(6.36), Inches(4.0), Inches(0.4),
      13, RED, bold=True, align=PP_ALIGN.CENTER)

    hint(sl, "SUGGERIMENTO: inserire diagramma di rete che mostra i confini della rete locale con firewall")
    return sl


def slide_17_architecture(prs):
    """Architettura integrata — UYUNI · n8n · AI."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), BLUE)

    T(sl, "L'Architettura Integrata — UYUNI · n8n · AI",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      26, WHITE, bold=True, align=PP_ALIGN.CENTER)

    comps = [
        ("UYUNI\n+\nSaltStack",  BLUE,   "Registra sistemi\nDistribuisce patch\nXML-RPC API"),
        ("XML-RPC\nConnector",   LGRAY,  "Ponte di\ncomunicazione\nper n8n"),
        ("n8n\nWorkflows",       TEAL,   "Orchestrazione AI\nTrigger · Audit\nWorkflow visuale"),
        ("AI Models\n(locale)",  GREEN,  "Analisi CVE\nRisk Score\nSuggerimenti"),
        ("Operatore\n+ Report",  ORANGE, "Approva azioni\nRiceve report\nDashboard"),
    ]
    bw, bh = Inches(2.25), Inches(3.1)
    gap    = Inches(0.22)
    sx, y  = Inches(0.35), Inches(1.4)

    for i, (label, color, desc) in enumerate(comps):
        x = sx + i * (bw + gap + Inches(0.15))
        R(sl, x, y, bw, bh, CARD, rx=True)
        R(sl, x, y, bw, Inches(0.12), color)
        T(sl, label, x, y + Inches(0.2), bw, Inches(1.0),
          14, color, bold=True, align=PP_ALIGN.CENTER)
        T(sl, desc,  x + Inches(0.12), y + Inches(1.3), bw - Inches(0.24), Inches(1.6),
          12, LGRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            ARR(sl, x + bw + Inches(0.08), y + bh / 2 - Inches(0.2),
                Inches(0.28), Inches(0.35), color)

    # Scenario descrittivo
    R(sl, Inches(0.5), Inches(4.75), Inches(12.3), Inches(2.05), CARD2, rx=True)
    T(sl, "Scenario tipo:",
      Inches(0.8), Inches(4.88), Inches(3), Inches(0.42),
      13, TEAL, bold=True)
    T(sl, ("UYUNI rileva nuova errata critica  →  XML-RPC notifica n8n  →  workflow AI analizza CVE + sistemi affetti\n"
           "→  modello locale calcola risk score  →  operatore riceve report strutturato con azione consigliata\n"
           "→  approva o il workflow esegue in autonomia  →  log completo dell'intera catena decisionale"),
      Inches(0.8), Inches(5.36), Inches(12.0), Inches(1.25),
      13, WHITE)

    hint(sl, "SUGGERIMENTO: sostituire con diagramma animato draw.io / Miro dell'architettura reale live")
    return sl


def slide_18_oss_ecosystem(prs):
    """Ecosistema Open Source — i numeri."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.08), YELLOW)

    T(sl, "L'Ecosistema Open Source — I Numeri Contano",
      Inches(0.6), Inches(0.18), Inches(12), Inches(0.65),
      26, WHITE, bold=True, align=PP_ALIGN.CENTER)

    tools = [
        ("n8n",        "50.000+", "stelle GitHub",
         "400+ integrazioni\n80.000+ community\n6 anni · Fair-code",       TEAL),
        ("SaltStack",  "14.000+", "stelle GitHub",
         "13+ anni in prod\nEvent-driven automation\nVMware backed",        ORANGE),
        ("PostgreSQL", "15.000+", "stelle GitHub",
         "35+ anni · Standard enterprise\nACID · Estensibile\nPGDG global", BLUE),
        ("UYUNI",      "800+",    "stelle GitHub",
         "15+ anni\nopenSUSE Foundation\nMulti-distro Linux",               GREEN),
        ("Flask",      "67.000+", "stelle GitHub",
         "15+ anni\nPallets Project\nMillioni di deployment",               YELLOW),
    ]
    tw, th = Inches(2.3), Inches(5.5)
    gap    = Inches(0.25)
    sx, y  = Inches(0.55), Inches(1.3)

    for i, (name, stars, label, desc, color) in enumerate(tools):
        x = sx + i * (tw + gap)
        R(sl, x, y, tw, th, CARD, rx=True)
        R(sl, x, y, tw, Inches(0.12), color)
        T(sl, name,  x, y + Inches(0.22), tw, Inches(0.55),
          17, color, bold=True, align=PP_ALIGN.CENTER)
        T(sl, stars, x, y + Inches(0.88), tw, Inches(0.95),
          34, WHITE, bold=True, align=PP_ALIGN.CENTER)
        T(sl, label, x + Inches(0.1), y + Inches(1.88), tw - Inches(0.2), Inches(0.42),
          12, LGRAY, align=PP_ALIGN.CENTER)
        R(sl, x + Inches(0.3), y + Inches(2.38), tw - Inches(0.6), Inches(0.04), CARD2)
        T(sl, desc,  x + Inches(0.15), y + Inches(2.55), tw - Inches(0.3), Inches(2.7),
          11, LGRAY, align=PP_ALIGN.CENTER)

    R(sl, Inches(0.5), H - Inches(0.55), Inches(12.3), Inches(0.42), CARD2, rx=True)
    T(sl, "Stabilità provata  ·  Community globale  ·  Nessun vendor lock-in  ·  Costo zero di licenza  ·  Codice auditabile",
      Inches(0.8), H - Inches(0.49), Inches(12.0), Inches(0.35),
      12, LGRAY, align=PP_ALIGN.CENTER)
    return sl


def slide_19_conclusions(prs):
    """Conclusioni."""
    sl = new_slide(prs)
    fill_bg(sl)
    R(sl, 0, 0, W, Inches(0.1), BLUE)
    R(sl, 0, H - Inches(0.1), W, Inches(0.1), TEAL)

    T(sl, "Conclusioni",
      Inches(0.6), Inches(0.22), Inches(12), Inches(0.75),
      34, WHITE, bold=True)

    rows = [
        (BLUE,   "①  INFRASTRUTTURA SOLIDA",
                 "UYUNI + SaltStack: base open source per gestione multi-OS, multi-tenant, CSP-independent"),
        (TEAL,   "②  AI CHE HA SENSO",
                 "Automazione, Intelligence, Novelty — ma solo con trasparenza e pieno controllo"),
        (ORANGE, "③  n8n: TRASPARENZA NATIVA",
                 "Ogni decisione AI è visibile, auditabile, debuggabile — workflow ingegnerizzato per segmento"),
        (GREEN,  "④  SICUREZZA PRIMA DI TUTTO",
                 "Modelli locali = zero esposizione dati clienti, GDPR, controllo totale sul modello"),
        (YELLOW, "⑤  OPEN SOURCE: LA SCELTA GIUSTA",
                 "50k+ stelle, community globale, stabilità enterprise, zero costi di licenza"),
    ]
    for i, (color, title, desc) in enumerate(rows):
        y = Inches(1.18) + i * Inches(1.12)
        R(sl, Inches(0.5), y, Inches(12.3), Inches(1.0), CARD, rx=True)
        R(sl, Inches(0.5), y, Inches(0.09), Inches(1.0), color)
        T(sl, title, Inches(0.75), y + Inches(0.05), Inches(4.8), Inches(0.46),
          14, color, bold=True)
        T(sl, desc,  Inches(0.75), y + Inches(0.55), Inches(11.8), Inches(0.38),
          13, LGRAY)

    T(sl, "Domande & Feedback",
      Inches(0.5), H - Inches(0.5), Inches(12.3), Inches(0.4),
      16, LGRAY, italic=True, align=PP_ALIGN.CENTER)
    return sl


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_title(prs)
    slide_02_benefits(prs)
    slide_03_reference_scenario(prs)
    slide_04_how_we_got_here(prs)
    slide_05_uyuni(prs)
    slide_06_oss_numbers(prs)
    slide_07_section_pillars(prs)
    slide_08_three_pillars(prs)
    slide_09_real_scenarios(prs)
    slide_10_orchestration(prs)
    slide_11_section_challenges(prs)
    slide_12_challenges(prs)
    slide_13_section_n8n(prs)
    slide_14_n8n_transparency(prs)
    slide_15_n8n_engineering(prs)
    slide_16_local_models(prs)
    slide_17_architecture(prs)
    slide_18_oss_ecosystem(prs)
    slide_19_conclusions(prs)

    import os
    out = os.path.join(os.path.dirname(__file__), "AI_IT_Operations_n8n.pptx")
    prs.save(out)
    print(f"✓  Salvato: {out}")
    print(f"   Slide totali: {len(prs.slides)}")


if __name__ == "__main__":
    main()
